import mesop as me
import mesop.labs as mel
import json
import asyncio
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, field
from enum import Enum
import httpx
import os
from io import BytesIO
import base64
from dotenv import load_dotenv
import openai

# PowerPoint generation imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load environment variables from .env
load_dotenv()

# Get AI Together API key from environment variable
ait_api_key = os.getenv("AIT_API_KEY")

if not ait_api_key:
    raise ValueError("AI Together API key is missing. Please set it in the .env file.")


# Data models
class ContentFormat(Enum):
    BULLETED_LIST = "bulleted_list"
    PARAGRAPH = "paragraph"


@dataclass
class SlideData:
    title: str
    content: str
    image_placeholder: str = "https://via.placeholder.com/400x300/cccccc/666666?text=Image+Placeholder"


@dataclass
class SlideTopic:
    title: str
    description: str = ""
    order: int = 0


@dataclass
class PresentationConfig:
    topic: str = ""
    num_slides: int = 5
    content_format: str = "bulleted_list"  # Store as string instead of enum
    background_color: str = "#ffffff"
    text_color: str = "#000000"
    audience: str = "colleagues"
    tone: str = "professional"
    scene: str = "general_scene"


@dataclass
class AppState:
    config: PresentationConfig = field(default_factory=PresentationConfig)
    slides: List[SlideData] = field(default_factory=list)
    slide_topics: List[SlideTopic] = field(default_factory=list)
    is_generating: bool = False
    is_generating_topics: bool = False
    error_message: str = ""
    generated_pptx: Optional[str] = None
    show_topic_breakdown: bool = False
    uploaded_file_content: str = ""
    download_filename: str = ""  # Add this
    download_ready: bool = False  # Add this

# Presentation options
AUDIENCE_OPTIONS = {
    "Superiors": "superiors",
    "Subordinates": "subordinates",
    "Colleagues": "colleagues",
    "Public": "public"
}

TONE_OPTIONS = {
    "Professional": "professional",
    "Friendly": "friendly",
    "Technical": "technical",
    "Persuasive": "persuasive",
    "Academic": "academic",
    "Inspirational": "inspirational",
    "Educational": "educational",
    "Humorous": "humorous",
    "Concise": "concise"
}

SCENE_OPTIONS = {
    "General Scene": "general_scene",
    "Teaching Materials": "teaching_materials",
    "Work Summary": "work_summary",
    "Work Plan": "work_plan",
    "Project Report": "project_report",
    "Solution": "solution",
    "Research Report": "research_report",
    "Meeting Materials": "meeting_materials",
    "Product Introduction": "product_introduction",
    "Company Introduction": "company_introduction",
    "Business Plan": "business_plan",
    "Science Popularization": "science_popularization",
    "Public Speaking": "public_speaking"
}

# Color options
BACKGROUND_COLORS = {
    "White": "#ffffff",
    "Light Gray": "#f5f5f5",
    "Dark Gray": "#333333",
    "Light Blue": "#e3f2fd",
    "Light Green": "#e8f5e9",
    "Light Purple": "#f3e5f5",
    "Light Orange": "#fff3e0",
    "Black": "#000000"
}

TEXT_COLORS = {
    "Black": "#000000",
    "Dark Gray": "#333333",
    "Blue": "#1976d2",
    "Green": "#388e3c",
    "Purple": "#7b1fa2",
    "Orange": "#f57c00",
    "Red": "#d32f2f",
    "White": "#ffffff"
}


# LLM Integration
class LLMClient:
    def __init__(self, api_key: str = None, base_url: str = "https://api.together.xyz/v1"):
        self.api_key = api_key or os.getenv("AIT_API_KEY", "your-api-key-here")
        self.base_url = base_url
        self.model = "meta-llama/Llama-3.3-70B-Instruct-Turbo"

    def generate_slide_topics(self, topic: str, num_slides: int, audience: str, tone: str, scene: str) -> List[
        SlideTopic]:
        """Generate slide topics breakdown using LLM API"""

        # Create context-aware prompt based on audience, tone, and scene
        audience_context = self._get_audience_context(audience)
        tone_context = self._get_tone_context(tone)
        scene_context = self._get_scene_context(scene)

        prompt = f"""
        Break down the following topic into {num_slides} slide topics for a presentation:

        Topic: "{topic}"

        Context:
        - Audience: {audience_context}
        - Tone: {tone_context}
        - Scene/Purpose: {scene_context}

        Requirements:
        - Create {num_slides} distinct slide topics that cover the main topic comprehensively
        - Each slide should have a clear, engaging title appropriate for the audience and tone
        - Include a brief description of what each slide should cover
        - Ensure logical flow from one slide to the next
        - Tailor the content structure to match the specified scene/purpose
        - Consider the audience level and adjust complexity accordingly
        - Apply the specified tone throughout the presentation structure

        Return the response as a JSON array with this structure:
        [
            {{
                "title": "Slide Topic Title",
                "description": "Brief description of what this slide covers",
                "order": 1
            }}
        ]

        Only return the JSON array, no additional text.
        """

        try:
            client = openai.OpenAI(
                api_key=self.api_key,
                base_url=self.base_url
            )
            response = client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            topics_data = json.loads(response.choices[0].message.content)
            return [
                SlideTopic(
                    title=topic_data["title"],
                    description=topic_data.get("description", ""),
                    order=topic_data.get("order", i + 1)
                )
                for i, topic_data in enumerate(topics_data)
            ]
        except Exception as e:
            # Fallback to mock data if API fails
            return self._generate_mock_topics(topic, num_slides)

    def _generate_mock_topics(self, topic: str, num_slides: int) -> List[SlideTopic]:
        """Generate mock slide topics for testing purposes"""
        mock_topics = []

        for i in range(num_slides):
            mock_topics.append(SlideTopic(
                title=f"{topic} - Topic {i + 1}",
                description=f"This slide will cover aspect {i + 1} of {topic}",
                order=i + 1
            ))

        return mock_topics

    def _get_audience_context(self, audience: str) -> str:
        """Get context description for audience type"""
        contexts = {
            "superiors": "Senior management, executives, or higher-level decision makers who prefer strategic overviews, clear outcomes, and executive summaries",
            "subordinates": "Direct reports, team members, or junior staff who need detailed guidance, clear instructions, and actionable steps",
            "colleagues": "Peers, team members at similar levels who understand the domain and appreciate collaborative discussion and technical details",
            "public": "General audience with varied backgrounds who need accessible language, clear explanations, and engaging content"
        }
        return contexts.get(audience, contexts["colleagues"])

    def _get_tone_context(self, tone: str) -> str:
        """Get context description for tone"""
        contexts = {
            "professional": "Formal, business-appropriate language with clear structure and respectful communication",
            "friendly": "Warm, approachable language that builds rapport while maintaining professionalism",
            "technical": "Precise, detailed language with industry-specific terminology and thorough explanations",
            "persuasive": "Compelling, convincing language that motivates action and builds strong arguments",
            "academic": "Scholarly, research-oriented language with evidence-based content and formal structure",
            "inspirational": "Motivating, uplifting language that energizes and encourages the audience",
            "educational": "Clear, instructional language that facilitates learning and understanding",
            "humorous": "Light-hearted, engaging language with appropriate humor to maintain audience interest",
            "concise": "Brief, direct language that delivers maximum impact with minimal words"
        }
        return contexts.get(tone, contexts["professional"])

    def _get_scene_context(self, scene: str) -> str:
        """Get context description for scene/purpose"""
        contexts = {
            "general_scene": "General presentation suitable for various contexts and audiences",
            "teaching_materials": "Educational content designed for learning and instruction with clear explanations and examples",
            "work_summary": "Professional summary of work completed, achievements, and outcomes for reporting purposes",
            "work_plan": "Strategic planning document outlining objectives, timelines, and action items for future work",
            "project_report": "Comprehensive project status update including progress, challenges, and next steps",
            "solution": "Problem-solving presentation that identifies issues and proposes actionable solutions",
            "research_report": "Academic or business research findings with data analysis and evidence-based conclusions",
            "meeting_materials": "Content designed for team meetings, discussions, and collaborative decision-making",
            "product_introduction": "Marketing and sales-focused content highlighting product features, benefits, and value proposition",
            "company_introduction": "Corporate presentation showcasing company overview, values, and capabilities",
            "business_plan": "Strategic business document outlining goals, strategies, market analysis, and financial projections",
            "science_popularization": "Educational content that makes complex scientific concepts accessible to general audiences",
            "public_speaking": "Engaging presentation designed for public venues with strong storytelling and audience interaction"
        }
        return contexts.get(scene, contexts["general_scene"])

    def generate_slides(self, slide_topics: List[SlideTopic], content_format: str, audience: str, tone: str,
                              scene: str) -> List[SlideData]:
        """Generate slides from slide topics using LLM API"""
        format_instruction = (
            "bulleted list format with 3-5 key points per slide"
            if content_format == "bulleted_list"
            else "paragraph format with 2-3 sentences per slide"
        )

        # Get context for better content generation
        audience_context = self._get_audience_context(audience)
        tone_context = self._get_tone_context(tone)
        scene_context = self._get_scene_context(scene)

        slides_data = []

        for topic in slide_topics:
            prompt = f"""
            Create detailed content for a presentation slide with the following specifications:

            Slide Title: "{topic.title}"
            Description: "{topic.description}"

            Context:
            - Audience: {audience_context}
            - Tone: {tone_context}
            - Scene/Purpose: {scene_context}

            Requirements:
            - Content should be in {format_instruction}
            - Apply the specified tone throughout the content
            - Tailor language and complexity for the target audience
            - Ensure content aligns with the presentation scene/purpose
            - Keep content concise but informative
            - Focus on the key points related to this specific slide topic

            Return the response as a JSON object with this structure:
            {{
                "title": "Slide Title",
                "content": "Slide content here"
            }}

            Only return the JSON object, no additional text.
            """

            try:
                client = openai.OpenAI(
                    api_key=self.api_key,
                    base_url=self.base_url
                )
                response = client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "user", "content": prompt}
                    ]
                )
                slide_data = json.loads(response.choices[0].message.content)

                slides_data.append(SlideData(
                    title=slide_data["title"],
                    content=slide_data["content"]
                ))

            except Exception as e:
                # Fallback to mock data if API fails
                if content_format == "bulleted_list":
                    content = f"• Key point 1 about {topic.title}\n• Important aspect to consider\n• Benefits and advantages\n• Future implications"
                else:
                    content = f"This slide covers {topic.description}. We'll explore the key concepts and their practical applications in this context."

                slides_data.append(SlideData(
                    title=topic.title,
                    content=content
                ))

        return slides_data


# PowerPoint Export
class PPTXExporter:
    def __init__(self, config: PresentationConfig):
        self.config = config

    def create_presentation(self, slides: List[SlideData]) -> str:
        """Create PowerPoint presentation from slides and return as base64 string"""
        print("Creating presentation")
        # Try to load template from disc first
        try:
            template_path = "Title.pptx"
            if os.path.exists(template_path):
                prs = Presentation(template_path)
                print("Using template from disc: template.pptx")
            else:
                print("No template found, using blank presentation")
                prs = Presentation()
        except Exception as e:
            print(f"Error loading disc template: {e}, using blank presentation")
            prs = Presentation()

        # Set slide size (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Get existing slides from template
        existing_slides = list(prs.slides)

        for i, slide_data in enumerate(slides):
            if i < len(existing_slides):
                # Use existing slide from template
                slide = existing_slides[i]
                print(f"Updating existing slide {i + 1}")
            else:
                # Add new slide if we have more data than template slides
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                print(f"Adding new slide {i + 1}")

            # Only modify background if it's a new slide or you want to override template
            if i >= len(existing_slides) or self.config.background_color != "#ffffff":
                background = slide.background
                fill = background.fill
                fill.solid()
                bg_color = self._hex_to_rgb(self.config.background_color)
                fill.fore_color.rgb = RGBColor(*bg_color)

            # Update title if title placeholder exists
            if slide.shapes.title:
                title = slide.shapes.title
                title.text = slide_data.title
                title.text_frame.paragraphs[0].font.size = Pt(32)
                title.text_frame.paragraphs[0].font.bold = True
                text_color = self._hex_to_rgb(self.config.text_color)
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*text_color)

            # Find and update content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_placeholder = shape
                    break

            if content_placeholder:
                content_placeholder.text = slide_data.content
                # Format content text
                for paragraph in content_placeholder.text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    text_color = self._hex_to_rgb(self.config.text_color)
                    paragraph.font.color.rgb = RGBColor(*text_color)

        # If template has more slides than our data, remove excess slides
        if len(existing_slides) > len(slides):
            # Remove slides from the end (in reverse order to maintain indices)
            for i in range(len(existing_slides) - 1, len(slides) - 1, -1):
                slide_id = existing_slides[i].slide_id
                prs.part.drop_rel(prs.slides._sldIdLst[i].rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])
                print(f"Removed excess slide {i + 1}")

        # Save to bytes and convert to base64
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return base64.b64encode(pptx_buffer.getvalue()).decode('utf-8')

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


# Initialize global state
_state = AppState()


def get_state() -> AppState:
    return _state


# Event handlers
def on_topic_change(e: me.InputEvent):
    state = get_state()
    state.config.topic = e.value
    state.show_topic_breakdown = False
    state.slide_topics = []


def on_num_slides_change(e: me.InputEvent):
    state = get_state()
    try:
        state.config.num_slides = int(e.value)
    except ValueError:
        state.config.num_slides = 5
    state.show_topic_breakdown = False
    state.slide_topics = []


def on_content_format_change(e: me.RadioChangeEvent):
    state = get_state()
    state.config.content_format = e.value


def on_background_color_change(e: me.SelectSelectionChangeEvent):
    state = get_state()
    state.config.background_color = e.value


def on_text_color_change(e: me.SelectSelectionChangeEvent):
    state = get_state()
    state.config.text_color = e.value


def on_audience_change(e: me.SelectSelectionChangeEvent):
    state = get_state()
    state.config.audience = e.value
    state.show_topic_breakdown = False
    state.slide_topics = []


def on_tone_change(e: me.SelectSelectionChangeEvent):
    state = get_state()
    state.config.tone = e.value
    state.show_topic_breakdown = False
    state.slide_topics = []


def on_scene_change(e: me.SelectSelectionChangeEvent):
    state = get_state()
    state.config.scene = e.value
    state.show_topic_breakdown = False
    state.slide_topics = []


def on_file_upload(e: me.UploadEvent):
    state = get_state()
    try:
        # Read the uploaded file content
        file_content = e.file.getvalue().decode('utf-8')
        state.uploaded_file_content = file_content
        state.config.topic = file_content
        state.show_topic_breakdown = False
        state.slide_topics = []
        state.error_message = "✅ File uploaded successfully!"
    except Exception as ex:
        state.error_message = f"Error reading file: {str(ex)}"


def on_use_uploaded_content(e: me.ClickEvent):
    state = get_state()
    if state.uploaded_file_content:
        state.config.topic = state.uploaded_file_content
        state.show_topic_breakdown = False
        state.slide_topics = []


def on_generate_topics(e: me.ClickEvent):
    state = get_state()

    if not state.config.topic.strip():
        state.error_message = "Please enter a topic for your presentation."
        return

    state.is_generating_topics = True
    state.error_message = ""
    state.slide_topics = []
    state.show_topic_breakdown = False
    try:
        # llm_client = LLMClient()
        # topics = llm_client.generate_slide_topics(
        #     state.config.topic,
        #     state.config.num_slides,
        #     state.config.audience,
        #     state.config.tone,
        #     state.config.scene
        # )
        topics = [SlideTopic(title='Introduction to Risk Management in Banking',
                             description='''Overview of the importance of risk management in the banking sector, its role in ensuring financial stability, and the types of risks banks face, inclu
ding credit, market, operational, and liquidity risks''',
                             order=1),
                  SlideTopic(title='Types of Risks in Banking: Understanding the Landscape',
                             description='''In-depth exploration of credit, market, operational, and liquidity risks, 
their characteristics, and the potential impact on banking operations, including examples and case studies''',
                             order=2)]

        state.slide_topics = topics
        state.show_topic_breakdown = True
    except Exception as e:
        state.error_message = f"Error generating topics: {str(e)}"
    finally:
        state.is_generating_topics = False

def on_slide_topic_change(index: int, e: me.InputEvent):
    state = get_state()
    if 0 <= index < len(state.slide_topics):
        state.slide_topics[index].title = e.value


def on_slide_description_change(index: int, e: me.InputEvent):
    state = get_state()
    if 0 <= index < len(state.slide_topics):
        state.slide_topics[index].description = e.value


def on_generate_slides(e: me.ClickEvent):
    state = get_state()

    if not state.slide_topics:
        state.error_message = "Please generate slide topics first."
        return

    state.is_generating = True
    state.error_message = ""
    state.slides = []
    state.generated_pptx = None  # Reset previous PPTX
    state.download_ready = False  # Reset download state

    try:
        # Generate slides (using your existing logic)
        # For demo purposes, keeping your mock data
        state.slides = [
            SlideData(title='Introduction to Risk Management in Banking',
                      content='''• Key point 1 about Introduction to Risk Management in Banking\n• Important aspect to consider\n• Benefits and advantages\n• Future implications''',
                      image_placeholder='https://via.placeholder.com/400x300/cccccc/666666?text=Image+Placeholder'),
            SlideData(title='Types of Risks in Banking: Understanding the Landscape',
                      content='• Key point 1 about Types of Risks in Banking: Understanding the Landscape\n• Important aspect to consider\n• Benefits and advantages\n• Future implications',
                      image_placeholder='https://via.placeholder.com/400x300/cccccc/666666?text=Image+Placeholder')
        ]

        # Generate PPTX
        exporter = PPTXExporter(state.config)
        state.generated_pptx = exporter.create_presentation(state.slides)

        # Prepare download immediately
        if state.generated_pptx:
            topic_clean = state.config.topic[:30].replace(' ', '_').replace('/', '_').replace('\\', '_')
            state.download_filename = f"presentation_{topic_clean}.pptx"
            state.download_ready = True
            state.error_message = "✅ Presentation generated successfully!"

    except Exception as e:
        state.error_message = f"Error generating slides: {str(e)}"
    finally:
        state.is_generating = False




def on_download_pptx(e: me.ClickEvent):
    state = get_state()
    if state.generated_pptx:
        # Generate a clean filename
        topic_clean = state.config.topic[:30].replace(' ', '_').replace('/', '_').replace('\\', '_')
        filename = f"presentation_{topic_clean}.pptx"

        try:
            # Set the download data in state for the UI to use
            state.download_filename = filename
            state.download_ready = True
            state.error_message = f"✅ '{filename}' ready for download."
        except Exception as e:
            state.error_message = f"Error preparing download: {str(e)}"




# UI Components
def header():
    with me.box(style=me.Style(
            background="#1976d2",
            color="white",
            padding=me.Padding.all(20),
            margin=me.Margin(bottom=20)
    )):
        me.text("Enhanced Slide AI", style=me.Style(
            font_size=32,
            font_weight="bold"
        ))
        me.text("Generate professional presentations with AI-powered topic breakdown", style=me.Style(
            font_size=16,
            opacity=0.9
        ))


def configuration_panel():
    state = get_state()

    with me.box(style=me.Style(
            background="white",
            border_radius=8,
            padding=me.Padding.all(20),
            margin=me.Margin(bottom=20),
            box_shadow="0 2px 4px rgba(0,0,0,0.1)"
    )):
        me.text("Configuration", style=me.Style(
            font_size=20,
            font_weight="bold",
            margin=me.Margin(bottom=15)
        ))

        # Topic input - now as textarea
        me.text("Presentation Topic:", style=me.Style(
            font_weight="bold",
            margin=me.Margin(bottom=5)
        ))
        me.textarea(
            label="Enter your presentation topic or content details",
            value=state.config.topic,
            on_input=on_topic_change,
            rows=6,
            style=me.Style(
                width="100%",
                margin=me.Margin(bottom=15),
                font_family="Arial, sans-serif"
            )
        )

        # File upload option
        me.text("Or upload a text file:", style=me.Style(
            font_weight="bold",
            margin=me.Margin(bottom=5)
        ))
        me.uploader(
            label="Upload Text File",
            accepted_file_types=[".txt", ".md"],
            on_upload=on_file_upload,
            style=me.Style(margin=me.Margin(bottom=10))
        )

        # Number of slides
        me.input(
            label="Number of Slides",
            value=str(state.config.num_slides),
            type="number",
            on_input=on_num_slides_change,
            style=me.Style(width="200px", margin=me.Margin(bottom=15))
        )

        # Content format
        me.text("Content Format:", style=me.Style(
            font_weight="bold",
            margin=me.Margin(bottom=5)
        ))
        me.radio(
            options=[
                me.RadioOption(label="Bulleted List", value="bulleted_list"),
                me.RadioOption(label="Paragraph", value="paragraph")
            ],
            value=state.config.content_format,
            on_change=on_content_format_change,
            style=me.Style(margin=me.Margin(bottom=15))
        )

        # Colors
        with me.box(style=me.Style(
                display="flex",
                gap=20,
                margin=me.Margin(bottom=20)
        )):
            me.select(
                label="Background Color",
                options=[
                    me.SelectOption(label=name, value=color)
                    for name, color in BACKGROUND_COLORS.items()
                ],
                value=state.config.background_color,
                on_selection_change=on_background_color_change,
                style=me.Style(width="200px")
            )

            me.select(
                label="Text Color",
                options=[
                    me.SelectOption(label=name, value=color)
                    for name, color in TEXT_COLORS.items()
                ],
                value=state.config.text_color,
                on_selection_change=on_text_color_change,
                style=me.Style(width="200px")
            )

        # Presentation context options
        me.text("Presentation Context:", style=me.Style(
            font_size=18,
            font_weight="bold",
            margin=me.Margin(top=20, bottom=15)
        ))

        with me.box(style=me.Style(
                display="flex",
                gap=20,
                margin=me.Margin(bottom=20)
        )):
            me.select(
                label="Target Audience",
                options=[
                    me.SelectOption(label=name, value=value)
                    for name, value in AUDIENCE_OPTIONS.items()
                ],
                value=state.config.audience,
                on_selection_change=on_audience_change,
                style=me.Style(width="200px")
            )

            me.select(
                label="Tone",
                options=[
                    me.SelectOption(label=name, value=value)
                    for name, value in TONE_OPTIONS.items()
                ],
                value=state.config.tone,
                on_selection_change=on_tone_change,
                style=me.Style(width="200px")
            )

            me.select(
                label="Scene/Purpose",
                options=[
                    me.SelectOption(label=name, value=value)
                    for name, value in SCENE_OPTIONS.items()
                ],
                value=state.config.scene,
                on_selection_change=on_scene_change,
                style=me.Style(width="200px")
            )

        # Helper text for context options
        with me.box(style=me.Style(
                background="#f0f8ff",
                border_radius=4,
                padding=me.Padding.all(12),
                margin=me.Margin(bottom=20)
        )):
            me.text("Context Help:", style=me.Style(
                font_weight="bold",
                font_size=14,
                margin=me.Margin(bottom=8)
            ))

            audience_desc = {
                "superiors": "For senior management - strategic focus, executive summaries",
                "subordinates": "For direct reports - detailed guidance, clear instructions",
                "colleagues": "For peers - collaborative discussion, technical details",
                "public": "For general audience - accessible language, engaging content"
            }

            tone_desc = {
                "professional": "Formal business language",
                "friendly": "Warm, approachable communication",
                "technical": "Precise, detailed explanations",
                "persuasive": "Compelling, convincing arguments",
                "academic": "Scholarly, research-oriented",
                "inspirational": "Motivating, uplifting",
                "educational": "Clear, instructional",
                "humorous": "Light-hearted, engaging",
                "concise": "Brief, direct communication"
            }

            scene_desc = {
                "general_scene": "General presentation",
                "teaching_materials": "Educational content",
                "work_summary": "Professional work reporting",
                "work_plan": "Strategic planning document",
                "project_report": "Project status update",
                "solution": "Problem-solving presentation",
                "research_report": "Research findings",
                "meeting_materials": "Team meeting content",
                "product_introduction": "Marketing presentation",
                "company_introduction": "Corporate overview",
                "business_plan": "Strategic business document",
                "science_popularization": "Science communication",
                "public_speaking": "Engaging public presentation"
            }

            me.text(f"• Audience: {audience_desc.get(state.config.audience, 'General audience')}", style=me.Style(
                font_size=12,
                margin=me.Margin(bottom=4)
            ))
            me.text(f"• Tone: {tone_desc.get(state.config.tone, 'Professional tone')}", style=me.Style(
                font_size=12,
                margin=me.Margin(bottom=4)
            ))
            me.text(f"• Scene: {scene_desc.get(state.config.scene, 'General presentation')}", style=me.Style(
                font_size=12
            ))

        # Generate topics button
        me.button(
            "Generate Slide Topics",
            on_click=on_generate_topics,
            disabled=state.is_generating_topics,
            style=me.Style(
                background="#ff9800",
                color="white",
                padding=me.Padding.symmetric(horizontal=20, vertical=10),
                border_radius=4,
                font_size=16,
                margin=me.Margin(right=10)
            )
        )

        if state.is_generating_topics:
            me.text("Generating slide topics...", style=me.Style(
                color="#ff9800",
                font_style="italic",
                margin=me.Margin(top=10)
            ))


def topic_breakdown_panel():
    state = get_state()

    if not state.show_topic_breakdown or not state.slide_topics:
        return

    with me.box(style=me.Style(
            background="white",
            border_radius=8,
            padding=me.Padding.all(20),
            margin=me.Margin(bottom=20),
            box_shadow="0 2px 4px rgba(0,0,0,0.1)"
    )):
        me.text("Slide Topics Breakdown", style=me.Style(
            font_size=20,
            font_weight="bold",
            margin=me.Margin(bottom=15)
        ))

        me.text("Review and edit the slide topics before generating the full presentation:", style=me.Style(
            color="#666666",
            margin=me.Margin(bottom=20)
        ))

        # Slide topics list
        for i, topic in enumerate(state.slide_topics):
            with me.box(style=me.Style(
                    border=me.Border.all(me.BorderSide(width=1, color="#e0e0e0", style="solid")),
                    border_radius=8,
                    padding=me.Padding.all(15),
                    margin=me.Margin(bottom=15),
                    background="#f9f9f9"
            )):
                me.text(f"Slide {i + 1}:", style=me.Style(
                    font_weight="bold",
                    margin=me.Margin(bottom=10)
                ))

                me.input(
                    label="Slide Title",
                    value=topic.title,
                    on_input=lambda e, idx=i: on_slide_topic_change(idx, e),
                    style=me.Style(width="100%", margin=me.Margin(bottom=10))
                )

                me.textarea(
                    label="Slide Description",
                    value=topic.description,
                    on_input=lambda e, idx=i: on_slide_description_change(idx, e),
                    rows=3,
                    style=me.Style(width="100%")
                )

        # Generate full slides button
        me.button(
            "Generate Full Presentation",
            on_click=on_generate_slides,
            disabled=state.is_generating,
            style=me.Style(
                background="#1976d2",
                color="white",
                padding=me.Padding.symmetric(horizontal=20, vertical=10),
                border_radius=4,
                font_size=16,
                margin=me.Margin(top=20)
            )
        )

        if state.is_generating:
            me.text("Generating full presentation...", style=me.Style(
                color="#1976d2",
                font_style="italic",
                margin=me.Margin(top=10)
            ))


def slides_preview():
    state = get_state()

    if not state.slides:
        return

    with me.box(style=me.Style(
            background="white",
            border_radius=8,
            padding=me.Padding.all(20),
            margin=me.Margin(bottom=20),
            box_shadow="0 2px 4px rgba(0,0,0,0.1)"
    )):
        with me.box(style=me.Style(
                display="flex",
                justify_content="space-between",
                align_items="center",
                margin=me.Margin(bottom=20)
        )):
            me.text("Slides Preview", style=me.Style(
                font_size=20,
                font_weight="bold"
            ))

            if state.generated_pptx:
                if state.download_ready:
                    # Create a data URL for direct download
                    data_url = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{state.generated_pptx}"

                    # Use me.html to create a direct download link
                    me.html(f'''
                                    <a href="{data_url}" 
                                       download="{state.download_filename}"
                                       style="
                                           background: #4caf50;
                                           color: white;
                                           padding: 8px 15px;
                                           border-radius: 4px;
                                           text-decoration: none;
                                           font-family: Arial, sans-serif;
                                           font-size: 14px;
                                       ">
                                       Download PPTX
                                    </a>
                                    ''')
                else:
                    me.button(
                        "Prepare Download",
                        on_click=on_download_pptx,
                        style=me.Style(
                            background="#4caf50",
                            color="white",
                            padding=me.Padding.symmetric(horizontal=15, vertical=8),
                            border_radius=4
                        )
                    )

        # PowerPoint Preview using iframe (if generated)
        if state.generated_pptx:
            # with me.box(style=me.Style(
            #         background="#f0f0f0",
            #         border_radius=8,
            #         padding=me.Padding.all(15),
            #         margin=me.Margin(bottom=20)
            # )):
            #     me.text("PowerPoint Preview", style=me.Style(
            #         font_size=16,
            #         font_weight="bold",
            #         margin=me.Margin(bottom=10)
            #     ))
            #
            #     # Create data URL for PowerPoint preview
            #     pptx_data_url = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{state.generated_pptx}"
            #
            #     # Use Office Online viewer for PowerPoint preview
            #     office_viewer_url = f"https://view.officeapps.live.com/op/embed.aspx?src={pptx_data_url}"
            #
            #     me.html(f'''
            #         <div style="width: 100%; height: 600px; border: 1px solid #ccc; border-radius: 8px; overflow: hidden;">
            #             <iframe
            #                 src="{office_viewer_url}"
            #                 width="100%"
            #                 height="100%"
            #                 frameborder="0"
            #                 style="border: none;">
            #                 Your browser does not support iframes. Please download the file to view it.
            #             </iframe>
            #         </div>
            #     ''')

            with me.box(style=me.Style(
                    background="#f0f0f0",
                    border_radius=8,
                    padding=me.Padding.all(15),
                    margin=me.Margin(bottom=20)
            )):
                me.text("PowerPoint Preview", style=me.Style(
                    font_size=16,
                    font_weight="bold",
                    margin=me.Margin(bottom=10)
                ))

                # Create a carousel-like preview
                me.html(f'''
                    <div style="
                        width: 100%; 
                        height: 600px; 
                        border: 2px solid #ccc; 
                        border-radius: 8px; 
                        background: white;
                        display: flex;
                        flex-direction: column;
                        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
                    ">
                        <div style="
                            background: #f8f9fa; 
                            padding: 10px; 
                            border-bottom: 1px solid #ccc;
                            font-family: Arial, sans-serif;
                            font-size: 14px;
                            color: #555;
                        ">
                            📊 PowerPoint Presentation ({len(state.slides)} slides) - Generated Successfully
                        </div>

                        <div style="
                            flex: 1;
                            display: flex;
                            align-items: center;
                            justify-content: center;
                            padding: 40px;
                            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                            color: white;
                            font-family: Arial, sans-serif;
                        ">
                        </div>

                        <div style="
                            background: #f8f9fa; 
                            padding: 15px; 
                            border-top: 1px solid #ccc;
                            font-family: Arial, sans-serif;
                            font-size: 12px;
                            color: #666;
                            text-align: center;
                        ">
                            💡 Your PowerPoint file is ready for download. The slides below show a text preview.
                        </div>
                    </div>
                ''')

        # Slides grid (text preview)
        me.text("Text Preview", style=me.Style(
            font_size=16,
            font_weight="bold",
            margin=me.Margin(top=20, bottom=15)
        ))

        for i, slide in enumerate(state.slides):
            with me.box(style=me.Style(
                    background=state.config.background_color,
                    color=state.config.text_color,
                    border_radius=8,
                    padding=me.Padding.all(20),
                    margin=me.Margin(bottom=15),
                    border=me.Border.all(me.BorderSide(width=1, color="#e0e0e0", style="solid")),
                    min_height="300px"
            )):
                me.text(f"Slide {i + 1}", style=me.Style(
                    font_size=12,
                    opacity=0.7,
                    margin=me.Margin(bottom=10)
                ))

                me.text(slide.title, style=me.Style(
                    font_size=24,
                    font_weight="bold",
                    margin=me.Margin(bottom=15)
                ))

                # Format content based on type
                if state.config.content_format == "bulleted_list":
                    for line in slide.content.split('\n'):
                        if line.strip():
                            me.text(line, style=me.Style(
                                font_size=16,
                                margin=me.Margin(bottom=5)
                            ))
                else:
                    me.text(slide.content, style=me.Style(
                        font_size=16,
                        line_height=1.5
                    ))


def error_display():
    state = get_state()

    if state.error_message:
        bg_color = "#d4edda" if state.error_message.startswith("✅") else "#f8d7da"
        text_color = "#155724" if state.error_message.startswith("✅") else "#721c24"
        border_color = "#c3e6cb" if state.error_message.startswith("✅") else "#f5c6cb"

        with me.box(style=me.Style(
                background=bg_color,
                color=text_color,
                border_radius=4,
                padding=me.Padding.all(15),
                margin=me.Margin(bottom=20),
                border=me.Border.all(me.BorderSide(width=1, color=border_color, style="solid"))
        )):
            me.text(state.error_message)


# Main app
@me.page(
    path="/",
    title="Enhanced Slide AI",
    security_policy=me.SecurityPolicy(
        allowed_iframe_parents=["https://google.github.io"]
    )
)
def main():
    with me.box(style=me.Style(
            background="#f5f5f5",
            min_height="100vh",
            font_family="Arial, sans-serif"
    )):
        header()

        with me.box(style=me.Style(
                max_width="1200px",
                margin=me.Margin.symmetric(horizontal="auto"),
                padding=me.Padding.symmetric(horizontal=20)
        )):
            error_display()
            configuration_panel()
            topic_breakdown_panel()
            slides_preview()


# if __name__ == "__main__":
#     me.run(main)